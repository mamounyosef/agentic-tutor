"""Content extraction tools for the Constructor agent system.

These tools allow the ingestion sub-agent to extract text from
various file types: PDFs, videos (via transcription), slides, documents.

Additionally provides tools for storing and organizing raw content files
in a structured course context folder hierarchy.
"""

import json
import os
import shutil
from pathlib import Path
from typing import Dict, Optional

from langchain_core.tools import tool

from app.core.config import get_settings
from app.core.transcription import transcribe_video


@tool
def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF file.

    Use this tool to process uploaded PDF course materials.
    Returns the full text content of the PDF.

    Args:
        file_path: Full path to the PDF file on disk

    Returns:
        JSON string with extracted text and metadata
    """
    try:
        import pypdf
    except ImportError:
        return json.dumps({
            "success": False,
            "text": "",
            "error": "pypdf not installed. Run: pip install pypdf"
        })

    try:
        if not os.path.exists(file_path):
            return json.dumps({
                "success": False,
                "text": "",
                "error": f"File not found: {file_path}"
            })

        reader = pypdf.PdfReader(file_path)
        page_count = len(reader.pages)

        # Extract text from all pages
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)

        full_text = "\n\n".join(text_parts)

        return json.dumps({
            "success": True,
            "text": full_text,
            "page_count": page_count,
            "file_path": file_path
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "text": "",
            "error": f"PDF extraction failed: {str(e)}"
        })


@tool
def extract_text_from_slides(file_path: str) -> str:
    """
    Extract text from PowerPoint slides (.ppt, .pptx).

    Use this tool to process uploaded presentation course materials.
    Returns the text content from all slides.

    Args:
        file_path: Full path to the presentation file on disk

    Returns:
        JSON string with extracted text and metadata
    """
    try:
        from pptx import Presentation
    except ImportError:
        return json.dumps({
            "success": False,
            "text": "",
            "error": "python-pptx not installed. Run: pip install python-pptx"
        })

    try:
        if not os.path.exists(file_path):
            return json.dumps({
                "success": False,
                "text": "",
                "error": f"File not found: {file_path}"
            })

        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext == ".ppt":
            return json.dumps({
                "success": False,
                "text": "",
                "error": (
                    "Legacy .ppt format is not supported. "
                    "Please convert the file to .pptx and re-upload it."
                )
            })

        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        # Extract text from all slides
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        full_text = "\n\n".join(text_parts)

        return json.dumps({
            "success": True,
            "text": full_text,
            "slide_count": slide_count,
            "file_path": file_path
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "text": "",
            "error": f"Slide extraction failed: {str(e)}"
        })


@tool
async def transcribe_video_file(file_path: str, language: Optional[str] = None) -> str:
    """
    Transcribe audio from a video file.

    Use this tool to extract spoken content from video course materials.
    Returns the full transcript text.

    Args:
        file_path: Full path to the video file on disk
        language: Optional language code (e.g., "en", "es", "auto")

    Returns:
        JSON string with transcription result and metadata
    """
    try:
        if not os.path.exists(file_path):
            return json.dumps({
                "success": False,
                "text": "",
                "error": f"File not found: {file_path}"
            })

        result = await transcribe_video(file_path, language)

        if result.get("error"):
            return json.dumps({
                "success": False,
                "text": result.get("text", ""),
                "error": result["error"]
            })

        return json.dumps({
            "success": True,
            "text": result.get("text", ""),
            "duration": result.get("duration", 0),
            "language": result.get("language"),
            "file_path": file_path
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "text": "",
            "error": f"Video transcription failed: {str(e)}"
        })


@tool
def extract_text_from_document(file_path: str) -> str:
    """
    Extract text from document files (.txt, .md, .docx).

    Use this tool to process text-based course materials.
    Returns the full text content.

    Args:
        file_path: Full path to the document file on disk

    Returns:
        JSON string with extracted text and metadata
    """
    try:
        if not os.path.exists(file_path):
            return json.dumps({
                "success": False,
                "text": "",
                "error": f"File not found: {file_path}"
            })

        file_ext = os.path.splitext(file_path)[1].lower()

        # Handle plain text files
        if file_ext in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

            return json.dumps({
                "success": True,
                "text": text,
                "file_type": file_ext,
                "file_path": file_path
            })

        # Handle Word documents
        elif file_ext == '.docx':
            try:
                from docx import Document
            except ImportError:
                return json.dumps({
                    "success": False,
                    "text": "",
                    "error": "python-docx not installed. Run: pip install python-docx"
                })

            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs if para.text])

            return json.dumps({
                "success": True,
                "text": text,
                "file_type": "docx",
                "file_path": file_path
            })

        else:
            return json.dumps({
                "success": False,
                "text": "",
                "error": f"Unsupported document type: {file_ext}"
            })

    except Exception as e:
        return json.dumps({
            "success": False,
            "text": "",
            "error": f"Document extraction failed: {str(e)}"
        })


@tool
def save_raw_content_to_file(
    course_id: int,
    original_filename: str,
    extracted_text: str,
    content_type: str = "other"
) -> str:
    """
    Save extracted raw content to a flat storage directory.

    This tool stores the extracted text in a temporary flat structure under
    /course_context_{course_id}/raw_content/. Later, the structure agent will
    organize these files into a module/unit hierarchy.

    Use this tool AFTER extracting text from any file (PDF, video, slides, document).

    Args:
        course_id: The ID of the course
        original_filename: Original name of the file (e.g., "intro_video.mp4")
        extracted_text: The full extracted text content to save
        content_type: Type of content - "pdf", "video", "slides", "document", or "other"

    Returns:
        JSON string with saved file path and status
    """
    try:
        # Resolve the course context directory from settings (always project root)
        course_context_dir = get_settings().course_context_absolute_path / f"course_context_{course_id}"
        raw_content_dir = course_context_dir / "raw_content"

        # Create directories if they don't exist
        raw_content_dir.mkdir(parents=True, exist_ok=True)

        # Determine output filename based on content type
        # For videos, append "_transcript" to make it clear
        base_name = Path(original_filename).stem  # Remove extension

        if content_type == "video":
            output_filename = f"{course_id}_{base_name}_transcript.txt"
        else:
            output_filename = f"{course_id}_{base_name}.txt"

        output_path = raw_content_dir / output_filename

        # Write the extracted text to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(extracted_text)

        return json.dumps({
            "success": True,
            "saved_path": str(output_path),
            "filename": output_filename,
            "content_type": content_type,
            "text_length": len(extracted_text),
            "message": f"Raw content saved to {output_filename}"
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to save raw content: {str(e)}"
        })


@tool
def organize_content_file(
    course_id: int,
    module_id: int,
    unit_id: int,
    original_filename: str,
    content_type: str = "other"
) -> str:
    """
    Organize a raw content file into the structured module/unit folder hierarchy.

    This tool moves a file from the flat /course_context_{course_id}/raw_content/
    directory into the organized structure:
    /course_context_{course_id}/module_{module_id}/unit_{unit_id}/

    The file is renamed to include all IDs for explicit tracking:
    {course_id}_{module_id}_{unit_id}_{original_filename}.txt

    IMPORTANT: Call this tool AFTER you have:
    1. Created the module (and received module_id)
    2. Created the unit (and received unit_id)
    3. Mapped this specific file to this unit in your blueprint

    Args:
        course_id: The ID of the course
        module_id: The ID of the module this file belongs to
        unit_id: The ID of the unit this file belongs to
        original_filename: Original name of the file (e.g., "intro_video.mp4")
        content_type: Type of content - "pdf", "video", "slides", "document", or "other"

    Returns:
        JSON string with organized file path and status
    """
    try:
        course_context_dir = get_settings().course_context_absolute_path / f"course_context_{course_id}"
        raw_content_dir = course_context_dir / "raw_content"

        # Determine source filename (must match what save_raw_content_to_file created)
        base_name = Path(original_filename).stem

        if content_type == "video":
            source_filename = f"{course_id}_{base_name}_transcript.txt"
        else:
            source_filename = f"{course_id}_{base_name}.txt"

        source_path = raw_content_dir / source_filename

        # Check if source file exists
        if not source_path.exists():
            return json.dumps({
                "success": False,
                "error": f"Source file not found: {source_path}. Make sure the file was saved during ingestion."
            })

        # Create the organized directory structure (simple IDs only)
        module_folder_name = f"module_{module_id}"
        unit_folder_name = f"unit_{unit_id}"

        organized_dir = course_context_dir / module_folder_name / unit_folder_name
        organized_dir.mkdir(parents=True, exist_ok=True)

        # Determine target filename with all IDs
        if content_type == "video":
            target_filename = f"{course_id}_{module_id}_{unit_id}_{base_name}_transcript.txt"
        else:
            target_filename = f"{course_id}_{module_id}_{unit_id}_{base_name}.txt"

        target_path = organized_dir / target_filename

        # Check if target already exists
        if target_path.exists():
            # File already organized - this is idempotent
            return json.dumps({
                "success": True,
                "organized_path": str(target_path),
                "filename": target_filename,
                "already_existed": True,
                "message": f"File already organized at {target_filename}"
            })

        # Move the file from flat structure to organized structure
        shutil.move(str(source_path), str(target_path))

        return json.dumps({
            "success": True,
            "organized_path": str(target_path),
            "filename": target_filename,
            "module_folder": module_folder_name,
            "unit_folder": unit_folder_name,
            "moved_from": str(source_path),
            "message": f"Content organized into {module_folder_name}/{unit_folder_name}/{target_filename}"
        })

    except Exception as e:
        return json.dumps({
            "success": False,
            "error": f"Failed to organize content file: {str(e)}"
        })
