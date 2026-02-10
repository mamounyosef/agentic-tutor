"""Ingestion tools for processing uploaded files.

Tools for parsing PDFs, PowerPoint presentations, Word documents, and videos.
"""

import hashlib
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

from langchain_core.tools import tool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)


class IngestionResult(BaseModel):
    """Result from an ingestion operation."""

    success: bool
    file_id: str
    file_type: str
    pages_or_slides: int
    content: str
    metadata: Dict[str, Any] = Field(default_factory=dict)
    errors: List[str] = Field(default_factory=list)


class ChunkResult(BaseModel):
    """Result from chunking content."""

    chunks: List[Dict[str, Any]]
    total_chunks: int
    avg_chunk_size: float


# =============================================================================
# PDF Ingestion
# =============================================================================

@tool
def ingest_pdf(file_path: str, course_id: int) -> Dict[str, Any]:
    """
    Parse a PDF file and extract text content.

    Args:
        file_path: Path to the PDF file
        course_id: ID of the course this file belongs to

    Returns:
        Dictionary with extracted content and metadata
    """
    try:
        # Lazy import to avoid dependency issues
        from pypdf import PdfReader
    except ImportError:
        return {
            "success": False,
            "error": "pypdf not installed. Run: pip install pypdf",
            "content": "",
            "metadata": {},
        }

    try:
        reader = PdfReader(file_path)
        pages = []

        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({
                "page_number": i + 1,
                "content": text,
            })

        # Combine all pages
        full_content = "\n\n".join(p["content"] for p in pages)

        # Generate file ID
        file_id = hashlib.md5(file_path.encode()).hexdigest()[:12]

        return {
            "success": True,
            "file_id": file_id,
            "file_type": "pdf",
            "pages_or_slides": len(pages),
            "content": full_content,
            "metadata": {
                "course_id": course_id,
                "original_filename": Path(file_path).name,
                "page_count": len(pages),
            },
        }

    except Exception as e:
        logger.error(f"Error ingesting PDF {file_path}: {e}")
        return {
            "success": False,
            "error": str(e),
            "content": "",
            "metadata": {},
        }


# =============================================================================
# PowerPoint Ingestion
# =============================================================================

@tool
def ingest_ppt(file_path: str, course_id: int) -> Dict[str, Any]:
    """
    Parse a PowerPoint file and extract slide content.

    Args:
        file_path: Path to the PPT/PPTX file
        course_id: ID of the course this file belongs to

    Returns:
        Dictionary with extracted content and metadata
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "success": False,
            "error": "python-pptx not installed. Run: pip install python-pptx",
            "content": "",
            "metadata": {},
        }

    try:
        prs = Presentation(file_path)
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            slides.append({
                "slide_number": i + 1,
                "content": "\n".join(slide_text),
            })

        # Combine all slides
        full_content = "\n\n--- SLIDE ---\n\n".join(
            f"Slide {s['slide_number']}:\n{s['content']}"
            for s in slides
        )

        file_id = hashlib.md5(file_path.encode()).hexdigest()[:12]

        return {
            "success": True,
            "file_id": file_id,
            "file_type": "pptx" if file_path.endswith(".pptx") else "ppt",
            "pages_or_slides": len(slides),
            "content": full_content,
            "metadata": {
                "course_id": course_id,
                "original_filename": Path(file_path).name,
                "slide_count": len(slides),
            },
        }

    except Exception as e:
        logger.error(f"Error ingesting PPT {file_path}: {e}")
        return {
            "success": False,
            "error": str(e),
            "content": "",
            "metadata": {},
        }


# =============================================================================
# Video Ingestion
# =============================================================================

@tool
async def ingest_video(
    file_path: str,
    course_id: int,
    transcript: Optional[str] = None,
    force_transcribe: bool = False,
) -> Dict[str, Any]:
    """
    Process a video file and extract full transcription.

    The video file is stored AS-IS for student viewing.
    The transcription is extracted for:
    - RAG searchability (find relevant content)
    - Quiz generation (questions from video content)
    - Student assessment (did they understand the video?)
    - Progress tracking (what did they learn from videos?)

    Transcription is performed using the configured service (OpenAI Whisper API by default).
    If a transcript is provided, it will be used directly unless force_transcribe=True.

    Args:
        file_path: Path to the video file (mp4, webm, mov, avi, etc.)
        course_id: ID of the course this file belongs to
        transcript: Optional pre-transcribed text (used if provided)
        force_transcribe: Force re-transcription even if transcript provided

    Returns:
        Dictionary with:
            - success: Whether ingestion succeeded
            - file_id: Unique identifier for the file
            - file_type: "video"
            - content: Full transcription text
            - metadata: Including duration, language, segments with timestamps
    """
    try:
        file_id = hashlib.md5(file_path.encode()).hexdigest()[:12]

        # Basic metadata
        metadata = {
            "course_id": course_id,
            "original_filename": Path(file_path).name,
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
            "file_path": file_path,
        }

        # If transcript is provided and we're not forcing re-transcription
        if transcript and not force_transcribe:
            logger.info(f"Using provided transcript for {file_path}")
            return {
                "success": True,
                "file_id": file_id,
                "file_type": "video",
                "pages_or_slides": 0,
                "content": transcript,
                "metadata": {
                    **metadata,
                    "transcript_source": "provided",
                },
            }

        # Perform transcription
        from ....core.transcription import transcribe_video

        logger.info(f"Transcribing video: {file_path}")
        transcription_result = await transcribe_video(file_path)

        # Check for errors
        if "error" in transcription_result and not transcription_result.get("text"):
            return {
                "success": False,
                "file_id": file_id,
                "file_type": "video",
                "pages_or_slides": 0,
                "content": "",
                "metadata": metadata,
                "error": transcription_result.get("error", "Transcription failed"),
                "warning": "Video transcription failed. Content not available for RAG.",
            }

        # Successful transcription
        transcript_text = transcription_result.get("text", "")
        metadata.update({
            "duration": transcription_result.get("duration", 0),
            "language": transcription_result.get("language"),
            "segments": transcription_result.get("segments", []),
            "transcript_source": transcription_result.get("service", "unknown"),
        })

        logger.info(
            f"Video transcription complete: {len(transcript_text)} chars, "
            f"{metadata['duration']}s, language={metadata['language']}"
        )

        return {
            "success": True,
            "file_id": file_id,
            "file_type": "video",
            "pages_or_slides": 0,
            "content": transcript_text,
            "metadata": metadata,
        }

    except Exception as e:
        logger.error(f"Error ingesting video {file_path}: {e}")
        return {
            "success": False,
            "file_id": file_id if "file_id" in locals() else "",
            "file_type": "video",
            "pages_or_slides": 0,
            "content": "",
            "metadata": {},
            "error": str(e),
        }


# =============================================================================
# DOCX Ingestion
# =============================================================================

@tool
def ingest_docx(file_path: str, course_id: int) -> Dict[str, Any]:
    """
    Parse a Word document and extract text content.

    Args:
        file_path: Path to the DOCX file
        course_id: ID of the course this file belongs to

    Returns:
        Dictionary with extracted content and metadata
    """
    try:
        from docx import Document
    except ImportError:
        return {
            "success": False,
            "error": "python-docx not installed. Run: pip install python-docx",
            "content": "",
            "metadata": {},
        }

    try:
        doc = Document(file_path)
        paragraphs = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Also extract tables
        tables_content = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            tables_content.append("\n".join(table_text))

        full_content = "\n\n".join(paragraphs)
        if tables_content:
            full_content += "\n\n--- TABLES ---\n\n" + "\n\n".join(tables_content)

        file_id = hashlib.md5(file_path.encode()).hexdigest()[:12]

        return {
            "success": True,
            "file_id": file_id,
            "file_type": "docx",
            "pages_or_slides": len(paragraphs),
            "content": full_content,
            "metadata": {
                "course_id": course_id,
                "original_filename": Path(file_path).name,
                "paragraph_count": len(paragraphs),
                "table_count": len(tables_content),
            },
        }

    except Exception as e:
        logger.error(f"Error ingesting DOCX {file_path}: {e}")
        return {
            "success": False,
            "error": str(e),
            "content": "",
            "metadata": {},
        }


# =============================================================================
# Text Chunking
# =============================================================================

@tool
def chunk_content_by_size(
    content: str,
    chunk_size: int = 1000,
    overlap: int = 200,
) -> Dict[str, Any]:
    """
    Split content into chunks by character size with overlap.

    Args:
        content: Text content to chunk
        chunk_size: Maximum characters per chunk (default 1000)
        overlap: Number of overlapping characters between chunks (default 200)

    Returns:
        Dictionary with list of chunks
    """
    if not content:
        return {"chunks": [], "total_chunks": 0, "avg_chunk_size": 0}

    chunks = []
    start = 0
    chunk_index = 0

    while start < len(content):
        end = start + chunk_size

        # Try to end at a sentence boundary
        if end < len(content):
            # Look for sentence end within a reasonable range
            search_start = max(end - 100, start)
            search_text = content[search_start:end]

            # Find last sentence ending punctuation
            for punct in [". ", "! ", "? ", "\n\n"]:
                last_punct = search_text.rfind(punct)
                if last_punct != -1:
                    end = search_start + last_punct + 1
                    break

        chunk_text = content[start:end].strip()
        if chunk_text:
            chunks.append({
                "chunk_id": f"chunk_{chunk_index}",
                "text": chunk_text,
                "start_char": start,
                "end_char": end,
                "chunk_type": "size_based",
            })
            chunk_index += 1

        start = end - overlap if end < len(content) else end

    avg_size = sum(len(c["text"]) for c in chunks) / len(chunks) if chunks else 0

    return {
        "chunks": chunks,
        "total_chunks": len(chunks),
        "avg_chunk_size": avg_size,
    }


@tool
def chunk_content_by_semantic(
    content: str,
    min_chunk_size: int = 200,
    max_chunk_size: int = 1500,
) -> Dict[str, Any]:
    """
    Split content into chunks based on semantic boundaries (headings, paragraphs).

    Args:
        content: Text content to chunk
        min_chunk_size: Minimum characters per chunk (default 200)
        max_chunk_size: Maximum characters per chunk (default 1500)

    Returns:
        Dictionary with list of semantically-grouped chunks
    """
    if not content:
        return {"chunks": [], "total_chunks": 0, "avg_chunk_size": 0}

    # Split by double newlines (paragraph/section boundaries)
    paragraphs = content.split("\n\n")

    chunks = []
    current_chunk = []
    current_size = 0
    chunk_index = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        para_size = len(para)

        # Check if this paragraph starts a new section (heading detection)
        is_heading = (
            para.startswith("#") or
            (len(para) < 100 and not para.endswith(".")) or
            para.isupper()
        )

        # Start new chunk if:
        # 1. Current chunk is big enough and we hit a heading
        # 2. Adding this paragraph would exceed max size
        if (is_heading and current_size >= min_chunk_size) or \
           (current_size + para_size > max_chunk_size and current_size > 0):
            # Save current chunk
            chunk_text = "\n\n".join(current_chunk)
            chunks.append({
                "chunk_id": f"chunk_{chunk_index}",
                "text": chunk_text,
                "chunk_type": "semantic",
            })
            chunk_index += 1
            current_chunk = []
            current_size = 0

        current_chunk.append(para)
        current_size += para_size + 2  # +2 for newlines

    # Don't forget the last chunk
    if current_chunk:
        chunk_text = "\n\n".join(current_chunk)
        chunks.append({
            "chunk_id": f"chunk_{chunk_index}",
            "text": chunk_text,
            "chunk_type": "semantic",
        })

    avg_size = sum(len(c["text"]) for c in chunks) / len(chunks) if chunks else 0

    return {
        "chunks": chunks,
        "total_chunks": len(chunks),
        "avg_chunk_size": avg_size,
    }


# =============================================================================
# Embeddings
# =============================================================================

@tool
async def generate_embeddings_for_chunks(
    chunks: List[Dict[str, Any]],
    course_id: int,
) -> Dict[str, Any]:
    """
    Generate embeddings for content chunks and prepare for vector storage.

    Args:
        chunks: List of chunk dictionaries with 'text' field
        course_id: ID of the course

    Returns:
        Dictionary with embedding status
    """
    from ....vector.constructor_store import ConstructorVectorStore

    try:
        vector_store = ConstructorVectorStore(course_id)

        # Add chunks to vector store (embeddings are generated automatically)
        chunk_ids = await vector_store.add_content_chunks(
            [
                {
                    "text": chunk["text"],
                    "chunk_type": chunk.get("chunk_type", "content"),
                    "source_file": chunk.get("source_file", ""),
                }
                for chunk in chunks
            ]
        )

        return {
            "success": True,
            "chunk_ids": chunk_ids,
            "total_embedded": len(chunk_ids),
        }

    except Exception as e:
        logger.error(f"Error generating embeddings: {e}")
        return {
            "success": False,
            "error": str(e),
            "chunk_ids": [],
            "total_embedded": 0,
        }
